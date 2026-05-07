#!/usr/bin/env python3
"""
populate_ebr.py
Reads a JSON payload of EBR data, populates named shapes in EBR_Template.pptx
and embeds screenshot images, then writes the result to stdout as base64.

Template: EBR_WIP.pptx (40 slides, named shapes) — rename to EBR_Template.pptx in repo
Usage: echo '<json>' | python3 populate_ebr.py
Last updated: May 2026
"""

import sys
import json
import os
import base64
import tempfile
import io
import re
import urllib.request

# ─── SHAPE NAME → NOTION FIELD MAP ───────────────────────────────────────────
# (slide_number, shape_name) → notion_field
# CRITICAL: Notion field names are exact — do not change spacing or casing.

SHAPE_MAP = {
    # ── Slide 1: Cover ────────────────────────────────────────────────────────
    (1,  "Text 1"):          "EBR Date",
    (1,  "Text 2"):          "__CUSTOMER_COVER__",           # special: replace token in existing text

    # ── Slide 4: Quote ────────────────────────────────────────────────────────
    (4,  "Quote"):           "Customer Quote",
    (4,  "Name"):            "Quote Attribution",

    # ── Slide 5: Recap ────────────────────────────────────────────────────────
    (5,  "Text 23"):         "Recap Point 1",
    (5,  "Text 28"):         "Recap Point 2",
    (5,  "Text 33"):         "Recap Point 3",

    # ── Slide 8: Customer Company Updates ─────────────────────────────────────
    (8,  "Text 2"):          "Customer Company Updates (1)",

    # ── Slide 9: High Level Snapshot ─────────────────────────────────────────
    (9,  "Stat 1"):          "Monthly Active Users",
    (9,  "Stat 2"):          "Average Weekly Users",
    (9,  "Stat 3"):          "Busiest Module",

    # ── Slide 10: Transactions Snapshot ──────────────────────────────────────
    (10, "Text 8"):          "Dup Caught Ahead of Pay Run (#)- Transactions",
    (10, "Text 12"):         "Dup Caught Ahead of Pay Run (\u00a3) - Transactions",
    (10, "Text 16"):         "Invoice Errors Ahead of Pay Run - Transactions",
    (10, "Text 20"):         "Most Common Error Type - Transactions",
    (10, "Text 24"):         "Most Active User - Transactions",

    # ── Slide 11: Helpdesk Snapshot ───────────────────────────────────────────
    (11, "Text 4"):          "Avg Handling Time - Current - Helpdesk",
    (11, "Text 8"):          "% Handling Time Decrease - Helpdesk",
    (11, "Text 12"):         "% Tickets via Gen AI - Helpdesk",
    (11, "Text 16"):         "% Tickets via Triggers - Helpdesk",
    (11, "Text 24"):         "Most Active User - Helpdesk",

    # ── Slide 12: Statements Snapshot ────────────────────────────────────────
    (12, "Text 4"):          "% Spend Reconciled - Statements",
    (12, "Text 8"):          "% Invoices Reconciled",
    (12, "Text 12"):         "% Vendors Reconciled - Statements",
    (12, "Text 16"):         "% Statements Reconciled ",         # trailing space — exact
    (12, "Text 20"):         "Missing Credits Recovered (\u00a3) - Statements",
    (12, "Text 24"):         "% Statements Automated",

    # ── Slide 35: Support Tickets ─────────────────────────────────────────────
    (35, "Text 7"):          "New Tickets Raised - Helpdesk",
    (35, "Text 18"):         "Open Tickets - Helpdesk",
    (35, "Text 20"):         "Tickets Waiting on Xelix - Helpdesk ",   # trailing space
    (35, "Text 26"):         "Tickets Waiting on Customer - Helpdesk",
    (35, "Text 35"):         "Closed Tickets - Helpdesk",

    # ── Slide 38: Renewal ─────────────────────────────────────────────────────
    (38, "Year 1TotalVal"):  "Contract Renewal Date",
    (38, "Year 2Pill"):      "Renewal Stakeholders",
}

# Shapes to clear
SHAPES_TO_CLEAR = [
    (5, "Text 24"),
    (5, "Text 29"),
    (5, "Text 34"),
]

# Summary slides — all get same combined block in Text 2
SUMMARY_SLIDES = [20, 27, 33]

# Screenshot field → (slide_number, position_index)
SCREENSHOT_MAP = {
    "Screenshot \u2013 Duplicate Invoices Chart":       (15, 0),
    "Screenshot \u2013 Duplicates by Cause":            (16, 0),
    "Screenshot \u2013 Duplicates Caught & Recovered":  (17, 0),
    "Screenshot \u2013 Invoice Errors Chart":           (18, 0),
    "Screenshot \u2013 Invoice Error Types":            (19, 0),
    "Screenshot \u2013 Invoice Errors by Cause":        (19, 1),
    "Screenshot \u2013 Reconciliations Chart":          (23, 0),
    "Screenshot \u2013 Industry Benchmark":             (24, 0),
    "Screenshot \u2013 Invoice Coverage":               (25, 0),
    "Screenshot \u2013 Full Automation Now":            (26, 0),
    "Screenshot \u2013 Full Automation Previous EBR 1": (26, 1),
    "Screenshot \u2013 Gen AI Outbound Chart":          (30, 0),
    "Screenshot \u2013 Avg Handling Time by Date":      (31, 0),
    "Screenshot \u2013 Avg Handling Time by User":      (31, 1),
    "Screenshot \u2013 Ticket Type Breakdown":          (32, 0),
    "Screenshot \u2013 Ticket Assignment Type":         (32, 1),
}


# ─── HELPERS ─────────────────────────────────────────────────────────────────

def get_val(data: dict, key: str, fallback: str = "") -> str:
    val = data.get(key, fallback)
    if val is None:
        return fallback
    if isinstance(val, list):
        return "\n".join(f"\u2022 {item}" for item in val if item)
    return str(val).strip()


def build_summary_block(data: dict) -> str:
    parts = []
    mapping = [
        ("Xelix Commitments",   "Xelix Commitments:"),
        ("Customer Commitments","Customer Commitments:"),
        ("Risks or Complaints", "Risks / Complaints:"),
        ("Recommended Action",  "Recommended Action:"),
    ]
    for key, label in mapping:
        val = get_val(data, key)
        if val:
            parts.append(f"{label}\n{val}")
    return "\n\n".join(parts)


def set_shape_text(shape, new_text: str):
    """
    Set text on a shape while preserving the font family (always Barlow),
    font size, bold and colour of the first existing run.
    Reduces to a single paragraph with a single run.
    """
    if not shape.has_text_frame:
        return

    from pptx.oxml.ns import qn
    from lxml import etree

    tf = shape.text_frame

    # Harvest formatting from first run
    font_name  = "Barlow"
    font_size  = None
    font_bold  = False
    font_color = "FFFFFF"

    for para in tf.paragraphs:
        for run in para.runs:
            if run.font.name and "Barlow" in run.font.name:
                font_name = run.font.name
            if run.font.size:
                font_size = run.font.size
            if run.font.bold is not None:
                font_bold = run.font.bold
            try:
                if run.font.color and run.font.color.type is not None:
                    font_color = str(run.font.color.rgb)
            except Exception:
                pass
            break
        break

    # Remove all paragraphs, keep one
    txBody = tf._txBody
    paras = txBody.findall(qn('a:p'))
    for p in paras[1:]:
        txBody.remove(p)

    para = tf.paragraphs[0]
    p_elem = para._p

    # Remove all runs from the paragraph
    for r in p_elem.findall(qn('a:r')):
        p_elem.remove(r)
    for br in p_elem.findall(qn('a:br')):
        p_elem.remove(br)

    # Build new run
    r_elem = etree.SubElement(p_elem, qn('a:r'))
    rpr = etree.SubElement(r_elem, qn('a:rPr'),
                           attrib={'lang': 'en-GB', 'dirty': '0'})

    if font_size:
        rpr.set('sz', str(int(font_size.pt * 100)))
    rpr.set('b', '1' if font_bold else '0')

    solidFill = etree.SubElement(rpr, qn('a:solidFill'))
    etree.SubElement(solidFill, qn('a:srgbClr'), attrib={'val': font_color})
    etree.SubElement(rpr, qn('a:latin'), attrib={'typeface': font_name})

    t_elem = etree.SubElement(r_elem, qn('a:t'))
    t_elem.text = new_text


# ─── POPULATION FUNCTIONS ────────────────────────────────────────────────────

def populate_shape_map(prs, data: dict):
    customer = get_val(data, "Customer")
    for (slide_idx, shape_name), notion_key in SHAPE_MAP.items():
        slide = prs.slides[slide_idx - 1]
        for shape in slide.shapes:
            if shape.name != shape_name or not shape.has_text_frame:
                continue
            if notion_key == "__CUSTOMER_COVER__":
                # Replace token in existing text (e.g. "Executive Business Review\nCustomer")
                existing = shape.text_frame.text
                new_text = existing.replace("Customer", customer) if customer else existing
                set_shape_text(shape, new_text)
            else:
                val = get_val(data, notion_key)
                if val:
                    set_shape_text(shape, val)
            break


def populate_cleared_shapes(prs):
    for (slide_idx, shape_name) in SHAPES_TO_CLEAR:
        slide = prs.slides[slide_idx - 1]
        for shape in slide.shapes:
            if shape.name == shape_name and shape.has_text_frame:
                set_shape_text(shape, "")
                break


def populate_summary_slides(prs, data: dict):
    summary = build_summary_block(data)
    if not summary:
        return
    for slide_idx in SUMMARY_SLIDES:
        slide = prs.slides[slide_idx - 1]
        for shape in slide.shapes:
            if shape.name == "Text 2" and shape.has_text_frame:
                set_shape_text(shape, summary)
                break


def populate_actions(prs, data: dict):
    slide = prs.slides[36]
    for i in range(1, 6):
        action = get_val(data, f"Action {i}")
        owner  = get_val(data, f"Action {i} Owner")
        due    = get_val(data, f"Action {i} Due")
        if not action:
            continue
        for shape in slide.shapes:
            if shape.name == f"Step {i}":
                set_shape_text(shape, action)
            if shape.name == f"Desc {i}":
                desc = " \u2014 ".join(filter(None, [owner, due]))
                set_shape_text(shape, desc if desc else action)


def populate_company_updates(prs, data: dict):
    slide = prs.slides[6]
    raw = get_val(data, "Xelix Company Updates")
    if not raw:
        return
    items = [s.strip() for s in re.split(r'[.\n]+', raw) if s.strip()][:6]
    for i, item in enumerate(items, 1):
        for shape in slide.shapes:
            if shape.name == f"Title {i}":
                set_shape_text(shape, f"0{i}  {item}")
            if shape.name == f"Desc {i}":
                set_shape_text(shape, "")


def download_image(url: str) -> bytes:
    req = urllib.request.Request(url, headers={'User-Agent': 'EBR-Generator/1.0'})
    with urllib.request.urlopen(req, timeout=15) as resp:
        return resp.read()


def embed_screenshots(prs, data: dict):
    for notion_field, (slide_idx, position) in SCREENSHOT_MAP.items():
        file_data = data.get(notion_field)
        if not file_data:
            continue

        # Resolve URL from various proxy formats
        url = None
        if isinstance(file_data, str) and file_data.startswith('http'):
            url = file_data
        elif isinstance(file_data, list) and len(file_data) > 0:
            item = file_data[0]
            if isinstance(item, dict):
                url = (item.get('url') or
                       item.get('file', {}).get('url') or
                       item.get('external', {}).get('url'))
            elif isinstance(item, str):
                url = item

        if not url:
            continue

        try:
            img_bytes = download_image(url)
        except Exception as e:
            print(f"Warning: could not download {notion_field}: {e}", file=sys.stderr)
            continue

        slide = prs.slides[slide_idx - 1]
        placeholders = [s for s in slide.shapes if s.name == "ImagePlaceholder"]

        if position >= len(placeholders):
            print(f"Warning: slide {slide_idx} has no ImagePlaceholder at index {position}", file=sys.stderr)
            continue

        ph = placeholders[position]
        left, top, width, height = ph.left, ph.top, ph.width, ph.height

        # Remove placeholder shape
        ph._element.getparent().remove(ph._element)

        # Embed image at same position
        slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, width, height)


# ─── MAIN ─────────────────────────────────────────────────────────────────────

def main():
    raw = sys.stdin.read().strip()
    if not raw:
        print("ERROR: no input", file=sys.stderr)
        sys.exit(1)

    try:
        data = json.loads(raw)
    except json.JSONDecodeError as e:
        print(f"ERROR: invalid JSON: {e}", file=sys.stderr)
        sys.exit(1)

    script_dir    = os.path.dirname(os.path.abspath(__file__))
    template_path = os.path.join(script_dir, "EBR_Template.pptx")
    if not os.path.exists(template_path):
        print(f"ERROR: EBR_Template.pptx not found at {template_path}", file=sys.stderr)
        sys.exit(1)

    from pptx import Presentation

    prs = Presentation(template_path)

    populate_shape_map(prs, data)
    populate_cleared_shapes(prs)
    populate_summary_slides(prs, data)
    populate_actions(prs, data)
    populate_company_updates(prs, data)
    embed_screenshots(prs, data)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        output_path = tmp.name

    try:
        prs.save(output_path)
        with open(output_path, 'rb') as f:
            b64 = base64.b64encode(f.read()).decode('utf-8')
        print(b64)
    finally:
        os.unlink(output_path)


if __name__ == "__main__":
    main()
